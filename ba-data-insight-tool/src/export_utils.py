from __future__ import annotations

from io import BytesIO

import pandas as pd
from plotly import graph_objects as go
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from .business_insights import insights_to_text


def export_excel(sheets: dict[str, pd.DataFrame], insights_text: str) -> bytes:
    output = BytesIO()
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        for name, df in sheets.items():
            safe_name = name[:31]
            if df is None or df.empty:
                pd.DataFrame({"mensaje": ["Sin datos disponibles"]}).to_excel(writer, sheet_name=safe_name, index=False)
            else:
                df.to_excel(writer, sheet_name=safe_name, index=False)
        pd.DataFrame({"insights": insights_text.splitlines()}).to_excel(writer, sheet_name="Insights ejecutivos", index=False)
        for sheet in writer.sheets.values():
            sheet.set_column(0, 20, 24)
    return output.getvalue()


def _fig_to_image(fig, width=640, height=320):
    """Convert a Plotly figure to PNG bytes via kaleido."""
    try:
        import plotly.io as pio
        return pio.to_image(fig, format="png", width=width, height=height, scale=1.5)
    except Exception:
        return None


def export_pdf(
    title: str,
    profile: dict,
    kpis: pd.DataFrame,
    insights: dict,
    figures: dict[str, go.Figure] | None = None,
) -> bytes:
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import HRFlowable, Image, PageBreak

    output = BytesIO()
    doc = SimpleDocTemplate(
        output,
        pagesize=letter,
        rightMargin=0.75 * inch,
        leftMargin=0.75 * inch,
        topMargin=0.75 * inch,
        bottomMargin=0.75 * inch,
    )
    styles = getSampleStyleSheet()
    NAVY = colors.HexColor("#1E3A5F")
    BLUE = colors.HexColor("#2E6DA4")
    LGRAY = colors.HexColor("#F2F5F8")

    title_style = ParagraphStyle("title_s", parent=styles["Title"], textColor=NAVY, fontSize=22, spaceAfter=6)
    h2_style = ParagraphStyle("h2_s", parent=styles["Heading2"], textColor=BLUE, fontSize=14, spaceBefore=14)
    body_style = ParagraphStyle("body_s", parent=styles["BodyText"], fontSize=10, leading=14)
    caption_style = ParagraphStyle("cap_s", parent=styles["BodyText"], fontSize=9, textColor=colors.grey, alignment=TA_CENTER)

    story = []

    # ── Portada ──────────────────────────────────────────────
    story.append(Spacer(1, 1 * inch))
    story.append(Paragraph(title, title_style))
    story.append(HRFlowable(width="100%", thickness=2, color=BLUE, spaceAfter=12))
    story.append(Paragraph("Reporte ejecutivo de análisis de datos", body_style))
    meta_data = [
        ["Registros analizados", str(profile.get("rows", "—"))],
        ["Columnas", str(profile.get("columns", "—"))],
        ["Datos faltantes", f"{profile.get('missing_percent', 0):.1f}%"],
        ["Duplicados", str(profile.get("duplicates", 0))],
    ]
    meta_table = Table(meta_data, colWidths=[2.5 * inch, 2.5 * inch])
    meta_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (0, -1), LGRAY),
        ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
        ("ROWBACKGROUNDS", (0, 0), (-1, -1), [colors.white, LGRAY]),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
    ]))
    story.extend([Spacer(1, 0.3 * inch), meta_table, PageBreak()])

    # ── KPIs ─────────────────────────────────────────────────
    story.append(Paragraph("KPIs principales", h2_style))
    story.append(HRFlowable(width="100%", thickness=1, color=LGRAY, spaceAfter=8))
    kpi_rows = [["KPI", "Valor", "Interpretación"]]
    for _, row in kpis.head(10).iterrows():
        kpi_rows.append([str(row.get("kpi", "—")), str(row.get("valor", "—")), str(row.get("interpretacion", "—"))])
    kpi_table = Table(kpi_rows, colWidths=[2.2 * inch, 1.4 * inch, 3.4 * inch])
    kpi_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), NAVY),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, LGRAY]),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    story.extend([kpi_table, PageBreak()])

    # ── Gráficos ─────────────────────────────────────────────
    if figures:
        story.append(Paragraph("Visualizaciones", h2_style))
        story.append(HRFlowable(width="100%", thickness=1, color=LGRAY, spaceAfter=8))
        for fig_name, fig in figures.items():
            img_bytes = _fig_to_image(fig)
            if img_bytes:
                img_io = BytesIO(img_bytes)
                img = Image(img_io, width=6.5 * inch, height=3.2 * inch)
                story.append(img)
                story.append(Paragraph(fig_name, caption_style))
                story.append(Spacer(1, 0.2 * inch))
        story.append(PageBreak())

    # ── Insights ─────────────────────────────────────────────
    story.append(Paragraph("Insights ejecutivos", h2_style))
    story.append(HRFlowable(width="100%", thickness=1, color=LGRAY, spaceAfter=8))
    for line in insights_to_text(insights).splitlines():
        if line.strip():
            style = h2_style if line.startswith("##") else body_style
            story.append(Paragraph(line.replace("##", "").strip(), style))
            story.append(Spacer(1, 4))

    doc.build(story)
    return output.getvalue()


def export_pptx(
    title: str,
    profile: dict,
    kpis: "pd.DataFrame",
    insights: dict,
    figures: "dict | None" = None,
    quality_score: "dict | None" = None,
) -> bytes:
    """Generate a professional PowerPoint presentation.

    Slides:
      1. Cover  2. KPIs  3+. Charts (one per figure)
      Last-1. Data quality  Last. Insights
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    NAVY  = RGBColor(0x1E, 0x3A, 0x5F)
    BLUE  = RGBColor(0x2E, 0x6D, 0xA4)
    TEAL  = RGBColor(0x1D, 0x9E, 0x75)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LGRAY = RGBColor(0xF2, 0xF5, 0xF8)
    BLACK = RGBColor(0x1A, 0x1A, 0x1A)
    RED   = RGBColor(0xC0, 0x39, 0x2B)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank layout

    def add_rect(slide, x, y, w, h, fill_color=None, line_color=None):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y),
                                        Inches(w), Inches(h))
        shape.line.fill.background()
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()
        return shape

    def add_text(slide, text, x, y, w, h, size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y),
                                          Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txBox

    def fig_to_image(fig, width=900, height=450):
        try:
            import plotly.io as pio
            return pio.to_image(fig, format="png",
                                width=width, height=height, scale=1.5)
        except Exception:
            return None

    # ── Slide 1: Cover ───────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=NAVY)
    add_rect(slide, 0, 5.8, 13.33, 1.7, fill_color=BLUE)
    add_text(slide, title, 0.8, 1.5, 11.5, 1.8,
             size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "Reporte ejecutivo de análisis de datos",
             0.8, 3.2, 11.5, 0.6, size=18, color=LGRAY)
    meta = (f"Registros: {profile.get('rows','—')}  ·  "
            f"Columnas: {profile.get('columns','—')}  ·  "
            f"Faltantes: {profile.get('missing_percent',0):.1f}%")
    add_text(slide, meta, 0.8, 6.0, 11.5, 0.6,
             size=13, color=WHITE)

    # ── Slide 2: KPIs ────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 1.0, fill_color=NAVY)
    add_text(slide, "KPIs principales", 0.4, 0.1, 12, 0.8,
             size=22, bold=True, color=WHITE)
    top_kpis = kpis.head(6)
    cols, col_w = 3, 4.0
    for i, (_, row) in enumerate(top_kpis.iterrows()):
        c, r = i % cols, i // cols
        x = 0.4 + c * col_w
        y = 1.3 + r * 2.8
        add_rect(slide, x, y, 3.8, 2.4, fill_color=LGRAY)
        add_rect(slide, x, y, 3.8, 0.08, fill_color=TEAL)
        add_text(slide, str(row.get("kpi","—")),
                 x+0.15, y+0.15, 3.5, 0.5,
                 size=11, color=BLUE, bold=True)
        add_text(slide, str(row.get("valor","—")),
                 x+0.15, y+0.65, 3.5, 0.8,
                 size=22, bold=True, color=NAVY)
        interp = str(row.get("interpretacion",""))[:80]
        add_text(slide, interp, x+0.15, y+1.5, 3.5, 0.8,
                 size=9, color=BLACK)

    # ── Slides 3+: Charts ────────────────────────────────────
    if figures:
        for fig_name, fig in figures.items():
            img_bytes = fig_to_image(fig)
            if not img_bytes:
                continue
            slide = prs.slides.add_slide(blank)
            add_rect(slide, 0, 0, 13.33, 0.9, fill_color=NAVY)
            add_text(slide, fig_name, 0.4, 0.08, 12, 0.75,
                     size=20, bold=True, color=WHITE)
            img_io = BytesIO(img_bytes)
            slide.shapes.add_picture(img_io, Inches(0.4), Inches(1.1),
                                      Inches(12.5), Inches(6.1))

    # ── Slide: Data quality ──────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 0.9, fill_color=NAVY)
    add_text(slide, "Calidad de datos", 0.4, 0.08, 12, 0.75,
             size=20, bold=True, color=WHITE)
    if quality_score:
        score = quality_score.get("score", 0)
        label = quality_score.get("label", "")
        score_color = (TEAL if score >= 80
                       else RGBColor(0xEF,0x9F,0x27) if score >= 55
                       else RED)
        add_rect(slide, 0.4, 1.1, 3.5, 2.5, fill_color=LGRAY)
        add_text(slide, f"{score}/100", 0.55, 1.4, 3.2, 1.2,
                 size=48, bold=True, color=score_color,
                 align=PP_ALIGN.CENTER)
        add_text(slide, label, 0.55, 2.6, 3.2, 0.5,
                 size=16, color=score_color, align=PP_ALIGN.CENTER)
    add_text(slide,
             f"Duplicados: {profile.get('duplicates',0)}  ·  "
             f"Faltantes globales: {profile.get('missing_percent',0):.1f}%",
             4.3, 1.2, 8.6, 0.6, size=14, color=BLACK)

    # ── Slide: Insights ──────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 0.9, fill_color=NAVY)
    add_text(slide, "Insights ejecutivos", 0.4, 0.08, 12, 0.75,
             size=20, bold=True, color=WHITE)
    add_rect(slide, 0, 0.9, 13.33, 6.6, fill_color=LGRAY)
    insight_text = insights_to_text(insights)
    lines = [l for l in insight_text.splitlines() if l.strip()][:14]
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.1),
                                      Inches(12.3), Inches(6.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        run = p.add_run()
        run.text = ("• " + line.replace("##","").strip()
                    if not line.startswith("##") else line.replace("##","").strip())
        run.font.size  = Pt(11 if not line.startswith("##") else 13)
        run.font.bold  = line.startswith("##")
        run.font.color.rgb = NAVY if line.startswith("##") else BLACK
        p.space_after  = Pt(4)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()
